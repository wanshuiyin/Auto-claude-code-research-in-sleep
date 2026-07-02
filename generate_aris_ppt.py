#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate a 3-slide high-density ARIS framework PPT in paper-talk style."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
import os

# -----------------------------------------------------------------------------
# Theme
# -----------------------------------------------------------------------------
PRIMARY = RgbColor(26, 35, 126)
ACCENT = RgbColor(255, 111, 97)
DARK = RgbColor(33, 33, 33)
GRAY = RgbColor(117, 117, 117)
PALE_BLUE = RgbColor(232, 234, 246)
WHITE = RgbColor(255, 255, 255)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.4)

# -----------------------------------------------------------------------------
# Helpers
# -----------------------------------------------------------------------------
def add_textbox(slide, left, top, width, height, text, font_size=12,
                bold=False, color=DARK, align=PP_ALIGN.LEFT,
                font_name="Calibri", italic=False, line_spacing=1.15):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    run.font.italic = italic
    return shape


def add_multiline(slide, left, top, width, height, lines, font_size=11,
                  color=DARK, line_spacing=1.15, bullet=False):
    """lines: list of strings."""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if bullet else "") + line
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return shape


def add_box(slide, left, top, width, height, title, body_lines,
            title_size=12, body_size=10, fill=PALE_BLUE, border=PRIMARY,
            title_color=PRIMARY, body_color=DARK):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = border
    box.line.width = Pt(1.2)
    add_textbox(slide, left + Inches(0.08), top + Inches(0.06),
                width - Inches(0.16), Inches(0.22),
                title, font_size=title_size, bold=True, color=title_color)
    add_multiline(slide, left + Inches(0.08), top + Inches(0.3),
                  width - Inches(0.16), height - Inches(0.38),
                  body_lines, font_size=body_size, color=body_color)
    return box


def add_arrow(slide, x1, y1, x2, y2, color=PRIMARY, width=Pt(1.2)):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = width
    line.line.end_arrowhead_width = 2
    line.line.end_arrowhead_length = 2
    return line


def add_footer(slide, page_num, total):
    add_textbox(slide, MARGIN, SLIDE_HEIGHT - Inches(0.28),
                Inches(2.5), Inches(0.2),
                "ARIS Framework & Design", font_size=9, color=GRAY)
    add_textbox(slide, SLIDE_WIDTH - Inches(0.8), SLIDE_HEIGHT - Inches(0.28),
                Inches(0.5), Inches(0.2),
                f"{page_num}/{total}", font_size=9, color=GRAY,
                align=PP_ALIGN.RIGHT)


def add_top_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                  SLIDE_WIDTH, Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()


def add_claim_title(slide, title):
    add_textbox(slide, MARGIN, Inches(0.15),
                SLIDE_WIDTH - 2 * MARGIN, Inches(0.5),
                title, font_size=22, bold=True, color=DARK)


def add_section_label(slide, left, top, text):
    add_textbox(slide, left, top, Inches(2.0), Inches(0.18),
                text.upper(), font_size=9, bold=True, color=ACCENT)

# -----------------------------------------------------------------------------
# Build
# -----------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]
    total = 3

    # =========================================================================
    # Slide 1: Design Philosophy
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    add_top_bar(slide)
    add_claim_title(slide, "ARIS Design Philosophy: Separate Progress from Judgment")

    # Central cross-model jury banner
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    MARGIN, Inches(0.78),
                                    SLIDE_WIDTH - 2 * MARGIN, Inches(0.42))
    banner.fill.solid(); banner.fill.fore_color.rgb = PRIMARY
    banner.line.fill.background()
    add_textbox(slide, MARGIN + Inches(0.1), Inches(0.84),
                SLIDE_WIDTH - 2 * MARGIN - Inches(0.2), Inches(0.3),
                "Cross-Model Jury: Claude executes → GPT-5.5 Codex adjudicates → deterministic scripts verify",
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Three columns
    col_w = (SLIDE_WIDTH - 2 * MARGIN - 2 * Inches(0.12)) / 3

    # Column 1: DRIVE vs ACQUIT
    add_section_label(slide, MARGIN, Inches(1.32), "Core rule")
    add_box(slide, MARGIN, Inches(1.5), col_w, Inches(2.55),
            "DRIVE vs ACQUIT",
            [
                "DRIVE: execute, schedule, generate, compile, mechanical checks",
                "Safe for the same model to perform",
                "ACQUIT: correctness, novelty, sufficiency, completeness",
                "Must be decided by a different model family",
                "One rule shapes every ARIS design decision"
            ],
            title_size=13, body_size=10, fill=PALE_BLUE)

    # Column 2: Gates
    add_section_label(slide, MARGIN + col_w + Inches(0.12), Inches(1.32), "Gates")
    add_box(slide, MARGIN + col_w + Inches(0.12), Inches(1.5), col_w, Inches(2.55),
            "Type-A / Type-B Gates",
            [
                "Type-A (mechanical): exit code, file exists, job completed",
                "Executor can self-judge safely",
                "Type-B (judgmental): quality, correctness, novelty",
                "Routed to GPT-5.5 Codex reviewer",
                "Shared 6-state verdict: PASS / WARN / FAIL / BLOCKED / ERROR / NOT_APPLICABLE"
            ],
            title_size=13, body_size=10, fill=WHITE, border=ACCENT, title_color=ACCENT)

    # Column 3: Axes
    add_section_label(slide, MARGIN + 2 * (col_w + Inches(0.12)), Inches(1.32), "Controls")
    ax_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     MARGIN + 2 * (col_w + Inches(0.12)), Inches(1.5),
                                     col_w, Inches(2.55))
    ax_box.fill.solid(); ax_box.fill.fore_color.rgb = PALE_BLUE
    ax_box.line.color.rgb = PRIMARY; ax_box.line.width = Pt(1.2)
    add_textbox(slide, MARGIN + 2 * (col_w + Inches(0.12)) + Inches(0.08), Inches(1.56),
                col_w - Inches(0.16), Inches(0.22),
                "Independent Axes", font_size=13, bold=True, color=PRIMARY)
    # effort row
    eff = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  MARGIN + 2 * (col_w + Inches(0.12)) + Inches(0.08), Inches(1.85),
                                  col_w - Inches(0.16), Inches(0.55))
    eff.fill.solid(); eff.fill.fore_color.rgb = WHITE
    eff.line.color.rgb = GRAY
    add_textbox(slide, MARGIN + 2 * (col_w + Inches(0.12)) + Inches(0.12), Inches(1.9),
                col_w - Inches(0.24), Inches(0.45),
                "effort\nlite → balanced → max → beast\n(how much work)",
                font_size=10, color=DARK)
    # assurance row
    ass = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  MARGIN + 2 * (col_w + Inches(0.12)) + Inches(0.08), Inches(2.5),
                                  col_w - Inches(0.16), Inches(0.55))
    ass.fill.solid(); ass.fill.fore_color.rgb = WHITE
    ass.line.color.rgb = GRAY
    add_textbox(slide, MARGIN + 2 * (col_w + Inches(0.12)) + Inches(0.12), Inches(2.55),
                col_w - Inches(0.24), Inches(0.45),
                "assurance\ndraft → polished → conference-ready → submission\n(how strict audits)",
                font_size=10, color=DARK)
    # example
    add_textbox(slide, MARGIN + 2 * (col_w + Inches(0.12)) + Inches(0.08), Inches(3.18),
                col_w - Inches(0.16), Inches(0.3),
                "Example: effort:lite, assurance:submission = fast run, strict gate",
                font_size=9, color=GRAY, italic=True)

    # Bottom relationship diagram
    rel_y = Inches(4.25)
    rel_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      MARGIN, rel_y,
                                      SLIDE_WIDTH - 2 * MARGIN, Inches(2.7))
    rel_box.fill.solid(); rel_box.fill.fore_color.rgb = WHITE
    rel_box.line.color.rgb = PRIMARY; rel_box.line.width = Pt(1.5)
    add_textbox(slide, MARGIN + Inches(0.1), rel_y + Inches(0.06),
                SLIDE_WIDTH - 2 * MARGIN - Inches(0.2), Inches(0.22),
                "How the ideas connect", font_size=12, bold=True, color=PRIMARY)

    # Four mini-boxes inside bottom area
    mini_w = (SLIDE_WIDTH - 2 * MARGIN - Inches(0.5)) / 4
    mini_h = Inches(1.9)
    mini_y = rel_y + Inches(0.35)
    mini_items = [
        ("1. No self-acquittal", ["Same model cannot judge", "its own output"]),
        ("2. Gate taxonomy", ["Mechanical vs judgmental", "determines who decides"]),
        ("3. Orthogonal controls", ["Effort and assurance", "scale independently"]),
        ("4. Verdict vocabulary", ["6 states shared across", "all skills and audits"])
    ]
    for i, (ttl, body) in enumerate(mini_items):
        x = MARGIN + Inches(0.1) + i * (mini_w + Inches(0.1))
        add_box(slide, x, mini_y, mini_w, mini_h, ttl, body,
                title_size=10, body_size=9, fill=PALE_BLUE)
        if i < len(mini_items) - 1:
            add_arrow(slide, x + mini_w + Inches(0.02), mini_y + mini_h / 2,
                      x + mini_w + Inches(0.08), mini_y + mini_h / 2,
                      color=GRAY, width=Pt(1.5))

    add_footer(slide, 1, total)

    # =========================================================================
    # Slide 2: Framework & Detailed Flow
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    add_top_bar(slide)
    add_claim_title(slide, "ARIS Framework & Detailed Flow: Six Workflows + Three Roles")

    # Left column: architecture + W1-W6
    left_w = Inches(5.0)
    add_section_label(slide, MARGIN, Inches(0.78), "Architecture")
    arch_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       MARGIN, Inches(0.95),
                                       left_w, Inches(2.3))
    arch_box.fill.solid(); arch_box.fill.fore_color.rgb = PALE_BLUE
    arch_box.line.color.rgb = PRIMARY; arch_box.line.width = Pt(1.5)

    # Three role boxes inside
    role_w = (left_w - Inches(0.4)) / 3
    role_h = Inches(0.75)
    role_y = Inches(1.15)
    roles = [
        ("Claude\nExecutor", "Write\nRun\nDraft"),
        ("GPT-5.5\nReviewer", "Review\nAudit\nAdjudicate"),
        ("Deterministic\nVerifier", "Exit-code\nGates\nScripts")
    ]
    for i, (ttl, body) in enumerate(roles):
        x = MARGIN + Inches(0.15) + i * (role_w + Inches(0.05))
        rb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, role_y, role_w, role_h)
        rb.fill.solid(); rb.fill.fore_color.rgb = WHITE
        rb.line.color.rgb = PRIMARY
        add_textbox(slide, x + Inches(0.03), role_y + Inches(0.04),
                    role_w - Inches(0.06), Inches(0.28),
                    ttl, font_size=8, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.03), role_y + Inches(0.34),
                    role_w - Inches(0.06), Inches(0.35),
                    body, font_size=8, color=DARK, align=PP_ALIGN.CENTER)

    # Arrow down to disk
    add_arrow(slide, MARGIN + left_w / 2, role_y + role_h,
              MARGIN + left_w / 2, Inches(2.2), color=PRIMARY)
    disk = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   MARGIN + Inches(0.3), Inches(2.2),
                                   left_w - Inches(0.6), Inches(0.45))
    disk.fill.solid(); disk.fill.fore_color.rgb = WHITE
    disk.line.color.rgb = ACCENT; disk.line.width = Pt(1.5)
    add_textbox(slide, MARGIN + Inches(0.35), Inches(2.27),
                left_w - Inches(0.7), Inches(0.35),
                "Disk state: receipts, REVIEW_STATE.json, traces", font_size=9,
                color=DARK, align=PP_ALIGN.CENTER)

    # W1-W6 pipeline below
    add_textbox(slide, MARGIN, Inches(2.78), left_w, Inches(0.2),
                "W1–W6 Lifecycle", font_size=11, bold=True, color=PRIMARY)

    stages = [
        ("W1", "Idea"),
        ("W1.5", "Exp"),
        ("W2", "Review"),
        ("S4", "Summary"),
        ("W3", "Paper"),
        ("W4", "Rebuttal"),
        ("W5", "Resubmit"),
        ("W6", "Talk")
    ]
    n = len(stages)
    gap = (left_w - Inches(0.2)) / n
    rad = Inches(0.22)
    y = Inches(3.05)
    for i, (code, name) in enumerate(stages):
        x = MARGIN + Inches(0.1) + i * gap + gap / 2 - rad
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, rad * 2, rad * 2)
        circ.fill.solid(); circ.fill.fore_color.rgb = PRIMARY if i < 4 else WHITE
        circ.line.color.rgb = PRIMARY
        add_textbox(slide, x, y + Inches(0.03), rad * 2, Inches(0.16),
                    code, font_size=7, bold=True,
                    color=WHITE if i < 4 else PRIMARY, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(0.2), rad * 2, Inches(0.22),
                    name, font_size=7, color=DARK, align=PP_ALIGN.CENTER)
        if i < n - 1:
            add_arrow(slide, x + rad * 2, y + rad,
                      MARGIN + Inches(0.1) + (i + 1) * gap + gap / 2 - rad, y + rad,
                      color=GRAY, width=Pt(0.8))

    # Outputs under pipeline
    outputs = ["IDEA_REPORT", "RESULTS", "AUTO_REVIEW", "NARRATIVE", "main.pdf", "PASTE_READY", "<Venue>", "slides"]
    for i, out in enumerate(outputs):
        x = MARGIN + Inches(0.1) + i * gap + gap / 2 - Inches(0.55)
        add_textbox(slide, x, y + rad * 2 + Inches(0.05), Inches(1.1), Inches(0.3),
                    out, font_size=6, color=GRAY, align=PP_ALIGN.CENTER)

    # Helper resolution chain at bottom left
    add_box(slide, MARGIN, Inches(3.95), left_w, Inches(1.05),
            "Integration Contract",
            [
                "6 components per skill integration: predicate + helper + artifact + checklist + backfill + verifier",
                "Helper resolution: .aris/tools → tools → $ARIS_REPO/tools → $CLAUDE_SKILL_DIR/scripts",
                "Failure policies A–E: block / warn-skip / forensic / cascade / diagnostic"
            ],
            title_size=10, body_size=9, fill=WHITE, border=ACCENT, title_color=ACCENT)

    # Right column: mechanisms
    right_x = MARGIN + left_w + Inches(0.2)
    right_w = SLIDE_WIDTH - right_x - MARGIN

    # Reviewer dispatch
    add_section_label(slide, right_x, Inches(0.78), "Reviewer dispatch")
    add_box(slide, right_x, Inches(0.95), right_w, Inches(1.05),
            "Cross-model review protocol",
            [
                "Reviewer prompt receives only file paths, never executor summaries",
                "Round 1: fresh Codex agent; Round 2+: continuation on same thread",
                "REVIEWER_BIAS_GUARD: improvement-loop uses fresh thread every round"
            ],
            title_size=10, body_size=9, fill=PALE_BLUE)

    # Fan-out
    add_section_label(slide, right_x, Inches(2.08), "Parallelism")
    add_box(slide, right_x, Inches(2.25), right_w, Inches(1.05),
            "Fan-out pattern",
            [
                "T1 Workflow parallel → T2 Agent-tool parallel → T3 Sequential fallback",
                "Shards EXTRACT candidates; mechanical dedup happens before jury",
                "Jury (Claude + GPT-5.5) ADJUDICATES, never the shards themselves"
            ],
            title_size=10, body_size=9, fill=WHITE, border=PRIMARY)

    # 5-layer audit
    add_section_label(slide, right_x, Inches(3.38), "Audit chain")
    audit_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        right_x, Inches(3.55),
                                        right_w, Inches(1.45))
    audit_box.fill.solid(); audit_box.fill.fore_color.rgb = PALE_BLUE
    audit_box.line.color.rgb = PRIMARY; audit_box.line.width = Pt(1.2)
    add_textbox(slide, right_x + Inches(0.08), Inches(3.61),
                right_w - Inches(0.16), Inches(0.2),
                "5-layer audit before submission", font_size=10, bold=True, color=PRIMARY)
    audits = [
        "1. experiment-audit",
        "2. result-to-claim",
        "3. paper-claim-audit",
        "4. citation-audit",
        "5. proof-checker / kill-argument"
    ]
    for i, txt in enumerate(audits):
        add_textbox(slide, right_x + Inches(0.08), Inches(3.85) + i * Inches(0.22),
                    right_w - Inches(0.16), Inches(0.2),
                    txt, font_size=9, color=DARK)
    # submission gate arrow
    gate = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   right_x + right_w - Inches(1.55), Inches(4.15),
                                   Inches(1.45), Inches(0.55))
    gate.fill.solid(); gate.fill.fore_color.rgb = ACCENT
    gate.line.fill.background()
    add_textbox(slide, right_x + right_w - Inches(1.5), Inches(4.22),
                Inches(1.35), Inches(0.45),
                "Submission\ngate exit 0", font_size=8, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_arrow(slide, right_x + right_w - Inches(1.65), Inches(4.42),
              right_x + right_w - Inches(1.6), Inches(4.42), color=ACCENT)

    add_footer(slide, 2, total)

    # =========================================================================
    # Slide 3: Memory System & Long-Cycle Operation
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    add_top_bar(slide)
    add_claim_title(slide, "ARIS Memory & Long-Cycle Operation: Survive, Resume, and Stay Honest")

    # Top banner: long-cycle principle
    top_ban = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      MARGIN, Inches(0.78),
                                      SLIDE_WIDTH - 2 * MARGIN, Inches(0.38))
    top_ban.fill.solid(); top_ban.fill.fore_color.rgb = PRIMARY
    top_ban.line.fill.background()
    add_textbox(slide, MARGIN + Inches(0.1), Inches(0.82),
                SLIDE_WIDTH - 2 * MARGIN - Inches(0.2), Inches(0.3),
                "Long-cycle principle: disk is the source of truth; timers fire-control but never judge",
                font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Three columns: Memory / Recovery / Operation
    col_w = (SLIDE_WIDTH - 2 * MARGIN - 2 * Inches(0.12)) / 3

    # Column 1: Memory system
    add_section_label(slide, MARGIN, Inches(1.25), "Memory system")
    add_box(slide, MARGIN, Inches(1.42), col_w, Inches(2.6),
            "Persistent knowledge & state",
            [
                "research_wiki.py: persistent store for ideas, claims, experiments, papers",
                "Wiki-helper-resolution: load/upsert variants with provenance authorization",
                "Capture antipatterns filter: no transient errors, no negative-capability claims",
                "State files per workflow:",
                "  • REVIEW_STATE.json (auto-review loop)",
                "  • PAPER_IMPROVEMENT_STATE.json (paper polish)",
                "  • REFINE_STATE.json (idea refine)",
                "  • queue_state.json + run_meta.txt (experiment queue)",
                "Receipts: .aris/runs/<run_id>.<phase>.done.json"
            ],
            title_size=11, body_size=9, fill=PALE_BLUE)

    # Column 2: Recovery
    add_section_label(slide, MARGIN + col_w + Inches(0.12), Inches(1.25), "Recovery")
    add_box(slide, MARGIN + col_w + Inches(0.12), Inches(1.42), col_w, Inches(2.6),
            "Resumable runs & cadence fence",
            [
                "run_state.py: separates done (executor finished) from accepted (gate passed)",
                "resume_point(): forward to first non-terminal phase",
                "Re-audit done-but-unaccepted stages on resume",
                "External cadence fence:",
                "  • /loop and CronCreate only fire-control progress",
                "  • Verdict-bearing skills are NEVER wrapped in /loop",
                "  • Heartbeat is Type-A: touch state, log iterations, nudge stalls",
                "Acceptance table maps each phase to its authority: codex agent-id or deterministic verifier path/sha"
            ],
            title_size=11, body_size=9, fill=WHITE, border=ACCENT, title_color=ACCENT)

    # Column 3: Operation
    add_section_label(slide, MARGIN + 2 * (col_w + Inches(0.12)), Inches(1.25), "Operation")
    add_box(slide, MARGIN + 2 * (col_w + Inches(0.12)), Inches(1.42), col_w, Inches(2.6),
            "Stall detection & orchestration",
            [
                "iteration_log.py: count new findings, not gut feelings",
                "Overnight heartbeat: self-target create_heartbeat",
                "  • stale≥2 → pivot=structural (change constraints)",
                "  • stale≥4 → pivot=human (escalate)",
                "watchdog.py: 24/7 health checks for remote training",
                "Paseo orchestration: create_agent → notifyOnFinish → gate → archive_agent",
                "Reviewer memory: continuation threads for loops, fresh threads for audits",
                "Review tracing: save_trace.sh writes 4 files + appends events.jsonl"
            ],
            title_size=11, body_size=9, fill=PALE_BLUE)

    # Bottom relationship diagram
    rel_y = Inches(4.2)
    rel_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      MARGIN, rel_y,
                                      SLIDE_WIDTH - 2 * MARGIN, Inches(2.75))
    rel_box.fill.solid(); rel_box.fill.fore_color.rgb = WHITE
    rel_box.line.color.rgb = PRIMARY; rel_box.line.width = Pt(1.5)
    add_textbox(slide, MARGIN + Inches(0.1), rel_y + Inches(0.06),
                SLIDE_WIDTH - 2 * MARGIN - Inches(0.2), Inches(0.22),
                "How memory, recovery, and orchestration reinforce each other", font_size=11,
                bold=True, color=PRIMARY)

    # Flow: Memory → State → Heartbeat → Paseo → Audit trail
    flow_y = rel_y + Inches(0.55)
    flow_h = Inches(0.7)
    flow_w = Inches(2.0)
    flow_items = [
        ("Memory layer", "wiki + state files"),
        ("Recovery layer", "resume + cadence fence"),
        ("Operation layer", "heartbeat + paseo"),
        ("Audit trail", "receipts + traces")
    ]
    gap2 = (SLIDE_WIDTH - 2 * MARGIN - 4 * flow_w - Inches(0.6)) / 3
    for i, (ttl, body) in enumerate(flow_items):
        x = MARGIN + Inches(0.15) + i * (flow_w + gap2 + Inches(0.2))
        fb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, flow_y, flow_w, flow_h)
        fb.fill.solid(); fb.fill.fore_color.rgb = PALE_BLUE
        fb.line.color.rgb = PRIMARY
        add_textbox(slide, x + Inches(0.05), flow_y + Inches(0.06),
                    flow_w - Inches(0.1), Inches(0.22),
                    ttl, font_size=10, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.05), flow_y + Inches(0.3),
                    flow_w - Inches(0.1), Inches(0.32),
                    body, font_size=8, color=DARK, align=PP_ALIGN.CENTER)
        if i < len(flow_items) - 1:
            add_arrow(slide, x + flow_w + Inches(0.03), flow_y + flow_h / 2,
                      x + flow_w + gap2 + Inches(0.17), flow_y + flow_h / 2,
                      color=ACCENT, width=Pt(2))

    # Reinforcement notes
    notes = [
        "Wiki prevents repeated ideas; state files enable crash recovery",
        "Cadence fence prevents timers from owning verdicts",
        "Paseo receipts make every agent action auditable",
        "Traces satisfy Policy C forensic requirements"
    ]
    note_y = rel_y + Inches(1.5)
    for i, note in enumerate(notes):
        col = i % 2
        row = i // 2
        x = MARGIN + Inches(0.15) + col * ((SLIDE_WIDTH - 2 * MARGIN - Inches(0.3)) / 2 + Inches(0.05))
        y = note_y + row * Inches(0.45)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.04), Inches(0.1), Inches(0.1))
        dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        add_textbox(slide, x + Inches(0.15), y,
                    (SLIDE_WIDTH - 2 * MARGIN - Inches(0.5)) / 2 - Inches(0.15), Inches(0.35),
                    note, font_size=9, color=DARK)

    add_footer(slide, 3, total)

    # Save
    os.makedirs("slides", exist_ok=True)
    os.makedirs(".aris/paper-talk", exist_ok=True)
    out_path = "slides/ARIS_Framework.pptx"
    prs.save(out_path)
    print(f"Saved {out_path} ({total} slides)")
    return out_path


if __name__ == "__main__":
    build()
